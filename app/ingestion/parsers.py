"""
Parsers for each supported document type.
Each parser returns a list of ParsedPage objects — one per page/slide.
The caller then runs vision passes and chunking on these.
"""
import io
import base64
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
import structlog

log = structlog.get_logger()


@dataclass
class ParsedPage:
    """
    Represents one page or slide extracted from a document.
    This is the unit that goes to the vision model and chunker.
    """
    page_number: int
    raw_text: str
    image_bytes: Optional[bytes] = None         # rendered image of the page/slide
    slide_title: Optional[str] = None           # PPT specific
    slide_notes: Optional[str] = None           # PPT speaker notes
    section_heading: Optional[str] = None       # PDF section heading detected above this page
    tables: list[str] = field(default_factory=list)         # tables as markdown strings
    code_blocks: list[str] = field(default_factory=list)    # for notebooks
    metadata: dict = field(default_factory=dict)


class PDFParser:
    """
    Layout-aware PDF parsing using PyMuPDF.
    - Extracts text with bounding boxes to preserve column structure
    - Detects section headings by font size heuristic
    - Extracts tables via pdfplumber
    - Renders each page as an image for the vision pass
    """

    def parse(self, file_bytes: bytes) -> list[ParsedPage]:
        import fitz  # PyMuPDF

        pages: list[ParsedPage] = []
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        current_section_heading = None

        for page_idx in range(len(doc)):
            page = doc[page_idx]

            # --- Text extraction with layout awareness ---
            blocks = page.get_text("dict")["blocks"]
            text_parts = []
            detected_heading = None

            for block in blocks:
                if block["type"] != 0:          # 0 = text block
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span["text"].strip()
                        if not text:
                            continue
                        font_size = span.get("size", 12)
                        # Heuristic: large bold text = section heading
                        flags = span.get("flags", 0)
                        is_bold = flags & 2 ** 4
                        if font_size >= 14 and is_bold and len(text) < 120:
                            detected_heading = text
                        text_parts.append(text)

            if detected_heading:
                current_section_heading = detected_heading

            raw_text = "\n".join(text_parts)

            # --- Render page as image for vision pass ---
            mat = fitz.Matrix(2.0, 2.0)         # 2x zoom = ~144 DPI, good for vision models
            pix = page.get_pixmap(matrix=mat)
            image_bytes = pix.tobytes("png")

            # --- Table extraction via pdfplumber ---
            tables_md = self._extract_tables_pdfplumber(file_bytes, page_idx)

            pages.append(ParsedPage(
                page_number=page_idx + 1,
                raw_text=raw_text,
                image_bytes=image_bytes,
                section_heading=current_section_heading,
                tables=tables_md,
                metadata={"page_width": page.rect.width, "page_height": page.rect.height}
            ))

            del pix  # explicit memory release

        doc.close()
        return pages

    def _extract_tables_pdfplumber(self, file_bytes: bytes, page_idx: int) -> list[str]:
        try:
            import pdfplumber
            tables_md = []
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                if page_idx >= len(pdf.pages):
                    return []
                page = pdf.pages[page_idx]
                for table in page.extract_tables():
                    if not table:
                        continue
                    md = self._table_to_markdown(table)
                    tables_md.append(md)
            return tables_md
        except Exception as e:
            log.warning("table_extraction_failed", page=page_idx, error=str(e))
            return []

    def _table_to_markdown(self, table: list[list]) -> str:
        if not table:
            return ""
        header = table[0]
        rows = table[1:]
        lines = []
        lines.append("| " + " | ".join(str(c or "") for c in header) + " |")
        lines.append("| " + " | ".join("---" for _ in header) + " |")
        for row in rows:
            lines.append("| " + " | ".join(str(c or "") for c in row) + " |")
        return "\n".join(lines)


class PPTXParser:
    """
    Parses PowerPoint files slide by slide.
    - Extracts title, bullet text, and speaker notes per slide
    - Renders each slide as an image (requires LibreOffice or python-pptx + pillow)
    - Treats each slide as a natural chunk boundary
    """

    def parse(self, file_bytes: bytes) -> list[ParsedPage]:
        from pptx import Presentation
        from pptx.util import Inches
        import fitz

        prs = Presentation(io.BytesIO(file_bytes))
        pages: list[ParsedPage] = []

        for slide_idx, slide in enumerate(prs.slides):
            title_text = ""
            body_parts = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                shape_text = shape.text_frame.text.strip()
                if not shape_text:
                    continue
                # Identify title shape
                if shape.shape_type == 13 or (hasattr(shape, "placeholder_format")
                        and shape.placeholder_format
                        and shape.placeholder_format.idx == 0):
                    title_text = shape_text
                else:
                    body_parts.append(shape_text)

            notes_text = ""
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    notes_text = notes_frame.text.strip()

            raw_text = "\n".join(filter(None, [title_text] + body_parts))

            # Render slide as image using LibreOffice (best quality)
            # Falls back to blank if LibreOffice not available
            image_bytes = self._render_slide_image(file_bytes, slide_idx)

            pages.append(ParsedPage(
                page_number=slide_idx + 1,
                raw_text=raw_text,
                image_bytes=image_bytes,
                slide_title=title_text or None,
                slide_notes=notes_text or None,
                metadata={"slide_index": slide_idx}
            ))

        return pages

    def _render_slide_image(self, file_bytes: bytes, slide_idx: int) -> Optional[bytes]:
        """
        Uses LibreOffice headless to convert PPTX to images.
        Requires: apt install libreoffice
        """
        import subprocess, tempfile, os
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                pptx_path = os.path.join(tmpdir, "deck.pptx")
                with open(pptx_path, "wb") as f:
                    f.write(file_bytes)
                subprocess.run([
                    "libreoffice", "--headless", "--convert-to", "png",
                    "--outdir", tmpdir, pptx_path
                ], capture_output=True, timeout=60)
                # LibreOffice generates one PNG per slide
                png_files = sorted(Path(tmpdir).glob("deck*.png"))
                if slide_idx < len(png_files):
                    return png_files[slide_idx].read_bytes()
        except Exception as e:
            log.warning("slide_render_failed", slide=slide_idx, error=str(e))
        return None


class DOCXParser:
    """Parses Word documents paragraph by paragraph, groups into logical pages."""

    def parse(self, file_bytes: bytes) -> list[ParsedPage]:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        all_paragraphs = []
        current_heading = None

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style.name.startswith("Heading"):
                current_heading = text
            all_paragraphs.append((text, current_heading))

        # Group ~30 paragraphs per "page" (simulated paging for DOCX)
        pages = []
        chunk_size = 30
        for i in range(0, len(all_paragraphs), chunk_size):
            group = all_paragraphs[i:i + chunk_size]
            raw_text = "\n".join(p for p, _ in group)
            section_heading = next((h for _, h in group if h), None)
            pages.append(ParsedPage(
                page_number=(i // chunk_size) + 1,
                raw_text=raw_text,
                section_heading=section_heading,
            ))

        return pages


class NotebookParser:
    """
    Parses Jupyter notebooks (.ipynb).
    Separates code cells from markdown cells.
    Treats each pair of markdown + code as a logical page.
    """

    def parse(self, file_bytes: bytes) -> list[ParsedPage]:
        import nbformat, json

        nb = nbformat.reads(file_bytes.decode("utf-8"), as_version=4)
        pages = []
        page_number = 0

        i = 0
        cells = nb.cells
        while i < len(cells):
            cell = cells[i]
            markdown_text = ""
            code_blocks = []

            # Collect consecutive markdown cells
            while i < len(cells) and cells[i].cell_type == "markdown":
                markdown_text += cells[i].source + "\n"
                i += 1

            # Collect following code cells
            while i < len(cells) and cells[i].cell_type == "code":
                code_blocks.append(cells[i].source)
                i += 1

            if not markdown_text and not code_blocks:
                i += 1
                continue

            page_number += 1
            raw_text = markdown_text.strip()

            pages.append(ParsedPage(
                page_number=page_number,
                raw_text=raw_text,
                code_blocks=code_blocks,
                metadata={"has_code": bool(code_blocks)}
            ))

        return pages


def get_parser(file_type: str):
    """Factory: returns the right parser for a given file type."""
    parsers = {
        "pdf":  PDFParser(),
        "pptx": PPTXParser(),
        "ppt":  PPTXParser(),
        "docx": DOCXParser(),
        "doc":  DOCXParser(),
        "ipynb": NotebookParser(),
    }
    parser = parsers.get(file_type.lower())
    if not parser:
        raise ValueError(f"Unsupported file type: {file_type}")
    return parser
